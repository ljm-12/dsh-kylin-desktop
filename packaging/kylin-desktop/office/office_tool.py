#!/usr/bin/env python3
"""Offline DOCX, XLSX, and PPTX inspection and editing utility."""

from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile
from copy import copy
from pathlib import Path
from typing import Any, Iterable


SUPPORTED = {".docx", ".xlsx", ".pptx", ".pdf"}


class OfficeToolError(RuntimeError):
    """A user-facing office operation failure."""


def emit(payload: dict[str, Any]) -> None:
    """Write one machine-readable result to stdout."""
    print(json.dumps(payload, ensure_ascii=False, indent=2))


def require_supported(path: Path) -> str:
    """Return a supported extension or raise a precise error."""
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED:
        raise OfficeToolError(f"unsupported file type: {suffix or '(none)'}; expected DOCX, XLSX, or PPTX")
    return suffix


def require_input(path: Path) -> str:
    """Validate an existing input file and return its extension."""
    suffix = require_supported(path)
    if not path.is_file():
        raise OfficeToolError(f"input file does not exist: {path}")
    return suffix


def require_distinct_output(source: Path, output: Path) -> None:
    """Require an explicit new output file to preserve the original."""
    if source.resolve() == output.resolve():
        raise OfficeToolError("output must differ from input; in-place editing is intentionally disabled")
    if source.suffix.lower() != output.suffix.lower():
        raise OfficeToolError("input and output extensions must match")


def atomic_save(save: Any, output: Path) -> None:
    """Save through a same-directory temporary file and atomically publish it."""
    output.parent.mkdir(parents=True, exist_ok=True)
    temp = output.with_name(f".{output.stem}.tmp-{os.getpid()}{output.suffix}")
    try:
        save(temp)
        os.replace(temp, output)
    finally:
        temp.unlink(missing_ok=True)


def find_matches(text: str, needle: str, replace_all: bool) -> list[tuple[int, int]]:
    """Return non-overlapping match spans."""
    if needle == "":
        raise OfficeToolError("find text must not be empty")
    matches: list[tuple[int, int]] = []
    start = 0
    while True:
        found = text.find(needle, start)
        if found < 0:
            return matches
        matches.append((found, found + len(needle)))
        if not replace_all:
            return matches
        start = found + len(needle)


def replace_in_runs(runs: Iterable[Any], needle: str, replacement: str, replace_all: bool = True) -> int:
    """Replace text across run boundaries while preserving surrounding run formatting."""
    run_list = list(runs)
    full_text = "".join(run.text or "" for run in run_list)
    matches = find_matches(full_text, needle, replace_all)
    if not matches:
        return 0

    spans: list[tuple[int, int]] = []
    cursor = 0
    for run in run_list:
        end = cursor + len(run.text or "")
        spans.append((cursor, end))
        cursor = end

    def locate(position: int, end_position: bool = False) -> tuple[int, int]:
        for index, (start, end) in enumerate(spans):
            if start <= position < end or (end_position and position == end and end > start):
                return index, position - start
        if position == len(full_text) and run_list:
            return len(run_list) - 1, len(run_list[-1].text or "")
        raise OfficeToolError("could not map an office text position to its formatting run")

    for start, end in reversed(matches):
        start_index, start_offset = locate(start)
        end_index, end_offset = locate(end, end_position=True)
        start_text = run_list[start_index].text or ""
        if start_index == end_index:
            run_list[start_index].text = start_text[:start_offset] + replacement + start_text[end_offset:]
            continue
        end_text = run_list[end_index].text or ""
        run_list[start_index].text = start_text[:start_offset] + replacement
        for index in range(start_index + 1, end_index):
            run_list[index].text = ""
        run_list[end_index].text = end_text[end_offset:]
    return len(matches)


def docx_paragraphs(document: Any) -> list[Any]:
    """Collect main, table, header, and footer paragraphs without duplicates."""
    paragraphs: list[Any] = []
    seen: set[int] = set()

    def add(paragraph: Any) -> None:
        identity = id(paragraph._p)
        if identity not in seen:
            seen.add(identity)
            paragraphs.append(paragraph)

    def add_table(table: Any) -> None:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    add(paragraph)
                for nested in cell.tables:
                    add_table(nested)

    for paragraph in document.paragraphs:
        add(paragraph)
    for table in document.tables:
        add_table(table)
    for section in document.sections:
        for part in (section.header, section.footer, section.first_page_header, section.first_page_footer):
            for paragraph in part.paragraphs:
                add(paragraph)
            for table in part.tables:
                add_table(table)
    return paragraphs


def pptx_text_paragraphs(presentation: Any) -> list[Any]:
    """Collect text paragraphs from shapes and tables."""
    paragraphs: list[Any] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                paragraphs.extend(shape.text_frame.paragraphs)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        paragraphs.extend(cell.text_frame.paragraphs)
    return paragraphs


def inspect_docx(path: Path) -> dict[str, Any]:
    """Return a bounded structural DOCX summary."""
    from docx import Document

    document = Document(path)
    paragraphs = document.paragraphs
    return {
        "type": "docx",
        "paragraphCount": len(paragraphs),
        "tableCount": len(document.tables),
        "sections": len(document.sections),
        "paragraphs": [
            {"index": index, "text": paragraph.text, "style": paragraph.style.name if paragraph.style else None}
            for index, paragraph in enumerate(paragraphs[:500])
        ],
        "tables": [
            {"index": index, "rows": len(table.rows), "columns": len(table.columns)}
            for index, table in enumerate(document.tables[:100])
        ],
    }


def inspect_xlsx(path: Path) -> dict[str, Any]:
    """Return workbook dimensions and a bounded non-empty cell sample."""
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=False, read_only=False)
    sheets: list[dict[str, Any]] = []
    for sheet in workbook.worksheets:
        cells: list[dict[str, Any]] = []
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value is not None:
                    cells.append({"cell": cell.coordinate, "value": cell.value, "numberFormat": cell.number_format})
                    if len(cells) >= 500:
                        break
            if len(cells) >= 500:
                break
        sheets.append(
            {
                "name": sheet.title,
                "maxRow": sheet.max_row,
                "maxColumn": sheet.max_column,
                "mergedRanges": [str(item) for item in sheet.merged_cells.ranges],
                "cells": cells,
            }
        )
    return {"type": "xlsx", "sheetCount": len(sheets), "sheets": sheets}


def inspect_pptx(path: Path) -> dict[str, Any]:
    """Return slide, shape, and visible text metadata."""
    from pptx import Presentation

    presentation = Presentation(path)
    slides: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(presentation.slides):
        shapes: list[dict[str, Any]] = []
        for shape_index, shape in enumerate(slide.shapes):
            item: dict[str, Any] = {
                "index": shape_index,
                "name": shape.name,
                "shapeType": str(shape.shape_type),
            }
            if getattr(shape, "has_text_frame", False):
                item["text"] = shape.text
            if getattr(shape, "has_table", False):
                item["table"] = {"rows": len(shape.table.rows), "columns": len(shape.table.columns)}
            shapes.append(item)
        slides.append({"index": slide_index, "shapes": shapes})
    return {"type": "pptx", "slideCount": len(slides), "slides": slides}


def inspect_pdf(path: Path) -> dict[str, Any]:
    """Inspect PDF page count and metadata."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        meta = {}
        if reader.metadata:
            for k, v in reader.metadata.items():
                if v:
                    meta[str(k).lstrip("/")] = str(v)
        return {"type": "pdf", "pageCount": len(reader.pages), "metadata": meta}
    except ImportError:
        return {"type": "pdf", "pageCount": "unknown", "note": "pypdf not installed"}
    except Exception as error:
        return {"type": "pdf", "error": str(error)}


def dump_docx(path: Path) -> str:
    """Extract DOCX paragraphs and tables as Markdown text."""
    from docx import Document

    document = Document(path)
    lines: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    for index, table in enumerate(document.tables):
        lines.append(f"\n### Table {index + 1}")
        for row_index, row in enumerate(table.rows):
            cells = [cell.text.strip().replace("\n", " ").replace("|", "\\|") for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
            if row_index == 0:
                lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n\n".join(lines)


def dump_xlsx(path: Path) -> str:
    """Extract XLSX sheet tables as Markdown text."""
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=True)
    sections: list[str] = []
    for sheet in workbook.worksheets:
        sections.append(f"# Sheet: {sheet.title}")
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        for row_index, row in enumerate(rows[:300]):
            cells = [str(c).replace("\n", " ").replace("|", "\\|") if c is not None else "" for c in row]
            if any(cells):
                sections.append("| " + " | ".join(cells) + " |")
                if row_index == 0:
                    sections.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n\n".join(sections)


def dump_pptx(path: Path) -> str:
    """Extract PPTX slides and tables as Markdown text."""
    from pptx import Presentation

    presentation = Presentation(path)
    sections: list[str] = []
    for slide_index, slide in enumerate(presentation.slides):
        sections.append(f"# Slide {slide_index + 1}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    sections.append(text)
            elif getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows):
                    cells = [cell.text.strip().replace("\n", " ").replace("|", "\\|") for cell in row.cells]
                    sections.append("| " + " | ".join(cells) + " |")
                    if row_index == 0:
                        sections.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n\n".join(sections)


def dump_pdf(path: Path) -> str:
    """Extract PDF pages as Markdown text."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(path)
        pages: list[str] = []
        for index, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"# Page {index + 1}\n\n{text.strip()}")
        return "\n\n".join(pages) if pages else "(No text extracted from PDF)"
    except ImportError:
        return "PDF text extraction requires pypdf library"
    except Exception as error:
        return f"PDF extraction error: {error}"


def dump_file(path: Path) -> dict[str, Any]:
    """Extract a supported office or PDF file as readable Markdown text."""
    suffix = require_input(path)
    if suffix == ".docx":
        content = dump_docx(path)
    elif suffix == ".xlsx":
        content = dump_xlsx(path)
    elif suffix == ".pptx":
        content = dump_pptx(path)
    elif suffix == ".pdf":
        content = dump_pdf(path)
    else:
        raise OfficeToolError(f"unsupported dump format: {suffix}")
    return {"ok": True, "path": str(path.resolve()), "type": suffix[1:], "content": content}


def inspect_file(path: Path) -> dict[str, Any]:
    """Inspect a supported office file."""
    suffix = require_input(path)
    if suffix == ".docx":
        result = inspect_docx(path)
    elif suffix == ".xlsx":
        result = inspect_xlsx(path)
    elif suffix == ".pptx":
        result = inspect_pptx(path)
    else:
        result = inspect_pdf(path)
    return {"ok": True, "path": str(path.resolve()), **result}


def replace_docx(source: Path, output: Path, needle: str, replacement: str, replace_all: bool) -> int:
    """Replace DOCX text in paragraphs, tables, headers, and footers."""
    from docx import Document

    document = Document(source)
    count = sum(replace_in_runs(paragraph.runs, needle, replacement, replace_all) for paragraph in docx_paragraphs(document))
    atomic_save(document.save, output)
    return count


def replace_xlsx(
    source: Path,
    output: Path,
    needle: str,
    replacement: str,
    replace_all: bool,
    include_formulas: bool,
) -> int:
    """Replace text in string cells while preserving workbook cell styles."""
    from openpyxl import load_workbook

    workbook = load_workbook(source, data_only=False, read_only=False)
    count = 0
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                value = cell.value
                if not isinstance(value, str) or (value.startswith("=") and not include_formulas):
                    continue
                matches = find_matches(value, needle, replace_all)
                if matches:
                    cell.value = value.replace(needle, replacement, -1 if replace_all else 1)
                    count += len(matches)
    atomic_save(workbook.save, output)
    return count


def replace_pptx(source: Path, output: Path, needle: str, replacement: str, replace_all: bool) -> int:
    """Replace PPTX text in shapes and tables while retaining run formatting."""
    from pptx import Presentation

    presentation = Presentation(source)
    count = sum(
        replace_in_runs(paragraph.runs, needle, replacement, replace_all)
        for paragraph in pptx_text_paragraphs(presentation)
    )
    atomic_save(presentation.save, output)
    return count


def apply_docx(source: Path, output: Path, operations: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Apply declarative DOCX operations."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    document = Document(source)
    results: list[dict[str, Any]] = []
    alignments = {
        "left": WD_ALIGN_PARAGRAPH.LEFT,
        "center": WD_ALIGN_PARAGRAPH.CENTER,
        "right": WD_ALIGN_PARAGRAPH.RIGHT,
        "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
    }

    def format_runs(paragraph: Any, spec: dict[str, Any]) -> None:
        for run in paragraph.runs:
            if "fontName" in spec:
                run.font.name = str(spec["fontName"])
            if "fontSizePt" in spec:
                run.font.size = Pt(float(spec["fontSizePt"]))
            for key in ("bold", "italic", "underline"):
                if key in spec:
                    setattr(run.font, key, bool(spec[key]))
            if "color" in spec:
                run.font.color.rgb = RGBColor.from_string(str(spec["color"]).lstrip("#"))
        if "alignment" in spec:
            value = str(spec["alignment"]).lower()
            if value not in alignments:
                raise OfficeToolError(f"unsupported paragraph alignment: {value}")
            paragraph.alignment = alignments[value]
        if "lineSpacing" in spec:
            paragraph.paragraph_format.line_spacing = float(spec["lineSpacing"])

    for op in operations:
        kind = op.get("op")
        if kind == "replace_text":
            count = sum(
                replace_in_runs(paragraph.runs, str(op["find"]), str(op.get("replace", "")), bool(op.get("all", True)))
                for paragraph in docx_paragraphs(document)
            )
            results.append({"op": kind, "replacements": count})
        elif kind == "set_paragraph":
            index = int(op["index"])
            paragraph = document.paragraphs[index]
            if paragraph.runs:
                paragraph.runs[0].text = str(op.get("text", ""))
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                paragraph.add_run(str(op.get("text", "")))
            if "style" in op:
                paragraph.style = str(op["style"])
            format_runs(paragraph, op.get("format", {}))
            results.append({"op": kind, "index": index})
        elif kind == "add_paragraph":
            paragraph = document.add_paragraph(str(op.get("text", "")), style=op.get("style"))
            format_runs(paragraph, op.get("format", {}))
            results.append({"op": kind, "index": len(document.paragraphs) - 1})
        elif kind == "set_table_cell":
            table_index = int(op["table"])
            row = int(op["row"])
            column = int(op["column"])
            document.tables[table_index].cell(row, column).text = str(op.get("text", ""))
            results.append({"op": kind, "table": table_index, "row": row, "column": column})
        elif kind == "format_paragraph":
            index = int(op["index"])
            format_runs(document.paragraphs[index], op.get("format", {}))
            results.append({"op": kind, "index": index})
        else:
            raise OfficeToolError(f"unsupported DOCX operation: {kind}")
    atomic_save(document.save, output)
    return results


def apply_xlsx(source: Path, output: Path, operations: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Apply declarative XLSX operations."""
    from openpyxl import load_workbook
    from openpyxl.styles import Alignment, PatternFill

    workbook = load_workbook(source, data_only=False, read_only=False)
    results: list[dict[str, Any]] = []

    for op in operations:
        kind = op.get("op")
        sheet = workbook[str(op.get("sheet", workbook.sheetnames[0]))]
        if kind == "set_cell":
            coordinate = str(op["cell"])
            cell = sheet[coordinate]
            if "formula" in op:
                formula = str(op["formula"])
                cell.value = formula if formula.startswith("=") else f"={formula}"
            elif "value" in op:
                cell.value = op["value"]
            style = op.get("style", {})
            if "numberFormat" in style:
                cell.number_format = str(style["numberFormat"])
            if "font" in style:
                font = copy(cell.font)
                for key, value in style["font"].items():
                    mapped = {"color": "color", "size": "sz", "name": "name"}.get(key, key)
                    setattr(font, mapped, value)
                cell.font = font
            if "fillColor" in style:
                color = str(style["fillColor"]).lstrip("#")
                cell.fill = PatternFill(fill_type="solid", fgColor=color)
            if "alignment" in style:
                alignment = copy(cell.alignment)
                for key, value in style["alignment"].items():
                    setattr(alignment, key, value)
                cell.alignment = alignment
            results.append({"op": kind, "sheet": sheet.title, "cell": coordinate})
        elif kind == "replace_text":
            count = 0
            needle = str(op["find"])
            replacement = str(op.get("replace", ""))
            replace_all = bool(op.get("all", True))
            for row in sheet.iter_rows():
                for cell in row:
                    value = cell.value
                    if not isinstance(value, str) or (value.startswith("=") and not op.get("includeFormulas", False)):
                        continue
                    matches = find_matches(value, needle, replace_all)
                    if matches:
                        cell.value = value.replace(needle, replacement, -1 if replace_all else 1)
                        count += len(matches)
            results.append({"op": kind, "sheet": sheet.title, "replacements": count})
        elif kind == "merge_cells":
            sheet.merge_cells(str(op["range"]))
            results.append({"op": kind, "sheet": sheet.title, "range": str(op["range"])})
        elif kind == "unmerge_cells":
            sheet.unmerge_cells(str(op["range"]))
            results.append({"op": kind, "sheet": sheet.title, "range": str(op["range"])})
        elif kind == "insert_rows":
            sheet.insert_rows(int(op["row"]), int(op.get("amount", 1)))
            results.append({"op": kind, "sheet": sheet.title})
        elif kind == "delete_rows":
            sheet.delete_rows(int(op["row"]), int(op.get("amount", 1)))
            results.append({"op": kind, "sheet": sheet.title})
        elif kind == "set_column_width":
            column = str(op["column"])
            sheet.column_dimensions[column].width = float(op["width"])
            results.append({"op": kind, "sheet": sheet.title, "column": column})
        elif kind == "set_row_height":
            row = int(op["row"])
            sheet.row_dimensions[row].height = float(op["height"])
            results.append({"op": kind, "sheet": sheet.title, "row": row})
        else:
            raise OfficeToolError(f"unsupported XLSX operation: {kind}")
    atomic_save(workbook.save, output)
    return results


def apply_pptx(source: Path, output: Path, operations: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Apply declarative PPTX operations."""
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.shared import Inches, Pt, RGBColor

    presentation = Presentation(source)
    results: list[dict[str, Any]] = []
    alignments = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}

    def locate_shape(op: dict[str, Any]) -> Any:
        slide = presentation.slides[int(op["slide"])]
        if "shapeName" in op:
            name = str(op["shapeName"])
            for shape in slide.shapes:
                if shape.name == name:
                    return shape
            raise OfficeToolError(f"shape not found: {name}")
        return slide.shapes[int(op["shape"])]

    def format_text_frame(text_frame: Any, spec: dict[str, Any]) -> None:
        for paragraph in text_frame.paragraphs:
            if "alignment" in spec:
                value = str(spec["alignment"]).lower()
                if value not in alignments:
                    raise OfficeToolError(f"unsupported PPTX alignment: {value}")
                paragraph.alignment = alignments[value]
            for run in paragraph.runs:
                if "fontName" in spec:
                    run.font.name = str(spec["fontName"])
                if "fontSizePt" in spec:
                    run.font.size = Pt(float(spec["fontSizePt"]))
                for key in ("bold", "italic", "underline"):
                    if key in spec:
                        setattr(run.font, key, bool(spec[key]))
                if "color" in spec:
                    run.font.color.rgb = RGBColor.from_string(str(spec["color"]).lstrip("#"))

    for op in operations:
        kind = op.get("op")
        if kind == "replace_text":
            count = sum(
                replace_in_runs(paragraph.runs, str(op["find"]), str(op.get("replace", "")), bool(op.get("all", True)))
                for paragraph in pptx_text_paragraphs(presentation)
            )
            results.append({"op": kind, "replacements": count})
        elif kind == "set_shape_text":
            shape = locate_shape(op)
            if not getattr(shape, "has_text_frame", False):
                raise OfficeToolError("selected shape has no text frame")
            shape.text = str(op.get("text", ""))
            format_text_frame(shape.text_frame, op.get("format", {}))
            results.append({"op": kind, "slide": int(op["slide"]), "shape": shape.name})
        elif kind == "format_shape_text":
            shape = locate_shape(op)
            if not getattr(shape, "has_text_frame", False):
                raise OfficeToolError("selected shape has no text frame")
            format_text_frame(shape.text_frame, op.get("format", {}))
            results.append({"op": kind, "slide": int(op["slide"]), "shape": shape.name})
        elif kind == "set_table_cell":
            shape = locate_shape(op)
            if not getattr(shape, "has_table", False):
                raise OfficeToolError("selected shape has no table")
            row = int(op["row"])
            column = int(op["column"])
            shape.table.cell(row, column).text = str(op.get("text", ""))
            results.append({"op": kind, "slide": int(op["slide"]), "row": row, "column": column})
        elif kind == "add_textbox":
            slide = presentation.slides[int(op["slide"])]
            shape = slide.shapes.add_textbox(
                Inches(float(op["left"])),
                Inches(float(op["top"])),
                Inches(float(op["width"])),
                Inches(float(op["height"])),
            )
            shape.text = str(op.get("text", ""))
            format_text_frame(shape.text_frame, op.get("format", {}))
            results.append({"op": kind, "slide": int(op["slide"]), "shape": shape.name})
        elif kind == "add_image":
            slide = presentation.slides[int(op["slide"])]
            image = Path(str(op["path"]))
            if not image.is_file():
                raise OfficeToolError(f"image does not exist: {image}")
            kwargs = {
                "left": Inches(float(op["left"])),
                "top": Inches(float(op["top"])),
            }
            if "width" in op:
                kwargs["width"] = Inches(float(op["width"]))
            if "height" in op:
                kwargs["height"] = Inches(float(op["height"]))
            shape = slide.shapes.add_picture(str(image), **kwargs)
            results.append({"op": kind, "slide": int(op["slide"]), "shape": shape.name})
        elif kind == "add_slide":
            layout = presentation.slide_layouts[int(op.get("layout", 1))]
            slide = presentation.slides.add_slide(layout)
            if "title" in op and slide.shapes.title is not None:
                slide.shapes.title.text = str(op["title"])
            body = str(op.get("body", ""))
            if body:
                for placeholder in slide.placeholders:
                    if placeholder != slide.shapes.title and getattr(placeholder, "has_text_frame", False):
                        placeholder.text = body
                        break
            results.append({"op": kind, "slide": len(presentation.slides) - 1})
        else:
            raise OfficeToolError(f"unsupported PPTX operation: {kind}")
    atomic_save(presentation.save, output)
    return results


def apply_plan(source: Path, output: Path, plan_path: Path) -> dict[str, Any]:
    """Load and execute a JSON office-edit plan."""
    suffix = require_input(source)
    require_distinct_output(source, output)
    try:
        plan = json.loads(plan_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as error:
        raise OfficeToolError(f"could not read plan JSON: {error}") from error
    operations = plan.get("operations")
    if not isinstance(operations, list) or not operations:
        raise OfficeToolError("plan must contain a non-empty operations array")
    if suffix == ".docx":
        results = apply_docx(source, output, operations)
    elif suffix == ".xlsx":
        results = apply_xlsx(source, output, operations)
    else:
        results = apply_pptx(source, output, operations)
    return {"ok": True, "input": str(source.resolve()), "output": str(output.resolve()), "operations": results}


def replace_file(args: argparse.Namespace) -> dict[str, Any]:
    """Dispatch a simple format-aware text replacement."""
    source = Path(args.input)
    output = Path(args.output)
    suffix = require_input(source)
    require_distinct_output(source, output)
    if suffix == ".docx":
        count = replace_docx(source, output, args.find, args.replace, not args.first)
    elif suffix == ".xlsx":
        count = replace_xlsx(source, output, args.find, args.replace, not args.first, args.include_formulas)
    else:
        count = replace_pptx(source, output, args.find, args.replace, not args.first)
    return {"ok": True, "input": str(source.resolve()), "output": str(output.resolve()), "replacements": count}


def create_file(args: argparse.Namespace) -> dict[str, Any]:
    """Create a simple new DOCX, XLSX, or PPTX document."""
    output = Path(args.output)
    suffix = require_supported(output)
    if output.exists():
        raise OfficeToolError(f"output already exists: {output}")
    output.parent.mkdir(parents=True, exist_ok=True)
    if suffix == ".docx":
        from docx import Document

        document = Document()
        if args.title:
            document.add_heading(args.title, level=0)
        if args.text:
            document.add_paragraph(args.text)
        atomic_save(document.save, output)
    elif suffix == ".xlsx":
        from openpyxl import Workbook

        workbook = Workbook()
        workbook.active.title = args.title or "Sheet1"
        if args.text:
            workbook.active["A1"] = args.text
        atomic_save(workbook.save, output)
    else:
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = args.title or ""
        for placeholder in slide.placeholders:
            if placeholder != slide.shapes.title and getattr(placeholder, "has_text_frame", False):
                placeholder.text = args.text or ""
                break
        atomic_save(presentation.save, output)
    return {"ok": True, "output": str(output.resolve()), "type": suffix[1:]}


def render_file(source: Path, output_dir: Path) -> dict[str, Any]:
    """Render through a locally installed LibreOffice with an isolated profile."""
    require_input(source)
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if executable is None:
        raise OfficeToolError("LibreOffice is not installed; rendering is unavailable on this deployment")
    output_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="dsh-office-lo-") as profile:
        command = [
            executable,
            "--headless",
            f"-env:UserInstallation=file://{profile}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(source),
        ]
        completed = subprocess.run(command, check=False, capture_output=True, text=True, timeout=180)
    expected = output_dir / f"{source.stem}.pdf"
    if completed.returncode != 0 or not expected.is_file():
        raise OfficeToolError(
            f"LibreOffice render failed with code {completed.returncode}: {(completed.stderr or completed.stdout).strip()}"
        )
    return {"ok": True, "input": str(source.resolve()), "output": str(expected.resolve())}


def build_parser() -> argparse.ArgumentParser:
    """Build the command-line parser."""
    parser = argparse.ArgumentParser(prog="dsh-office", description=__doc__)
    subparsers = parser.add_subparsers(dest="command", required=True)

    dump_parser = subparsers.add_parser("dump", help="dump document text and tables as Markdown")
    dump_parser.add_argument("input")

    inspect_parser = subparsers.add_parser("inspect", help="inspect document structure as JSON")
    inspect_parser.add_argument("input")

    replace_parser = subparsers.add_parser("replace", help="replace text and save a new file")
    replace_parser.add_argument("input")
    replace_parser.add_argument("output")
    replace_parser.add_argument("--find", required=True)
    replace_parser.add_argument("--replace", required=True)
    replace_parser.add_argument("--first", action="store_true")
    replace_parser.add_argument("--include-formulas", action="store_true")

    apply_parser = subparsers.add_parser("apply", help="apply a declarative JSON edit plan")
    apply_parser.add_argument("input")
    apply_parser.add_argument("output")
    apply_parser.add_argument("plan")

    create_parser = subparsers.add_parser("create", help="create a simple new office file")
    create_parser.add_argument("output")
    create_parser.add_argument("--title", default="")
    create_parser.add_argument("--text", default="")

    render_parser = subparsers.add_parser("render", help="render to PDF with installed LibreOffice")
    render_parser.add_argument("input")
    render_parser.add_argument("output_dir")
    return parser


def main() -> int:
    """Run one office command and report JSON."""
    parser = build_parser()
    args = parser.parse_args()
    try:
        if args.command == "dump":
            result = dump_file(Path(args.input))
        elif args.command == "inspect":
            result = inspect_file(Path(args.input))
        elif args.command == "replace":
            result = replace_file(args)
        elif args.command == "apply":
            result = apply_plan(Path(args.input), Path(args.output), Path(args.plan))
        elif args.command == "create":
            result = create_file(args)
        else:
            result = render_file(Path(args.input), Path(args.output_dir))
        emit(result)
        return 0
    except (OfficeToolError, IndexError, KeyError, ValueError) as error:
        emit({"ok": False, "error": str(error)})
        return 2


if __name__ == "__main__":
    sys.exit(main())
